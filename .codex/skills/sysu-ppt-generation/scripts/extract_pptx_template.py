#!/usr/bin/env python
"""Extract reusable design signals from local PowerPoint templates."""

from __future__ import annotations

import argparse
import json
import re
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation


EMU_PER_INCH = 914400
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def emu_to_in(value: int | None) -> float | None:
    if value is None:
        return None
    return round(value / EMU_PER_INCH, 3)


def relpath(path: Path, root: Path) -> str:
    try:
        return str(path.resolve().relative_to(root.resolve())).replace("\\", "/")
    except ValueError:
        return str(path.resolve()).replace("\\", "/")


def clean_text(text: str, limit: int = 80) -> str:
    value = re.sub(r"\s+", " ", text or "").strip()
    return value[: limit - 1] + "…" if len(value) > limit else value


def xml_root(zf: zipfile.ZipFile, name: str) -> etree._Element | None:
    try:
        return etree.fromstring(zf.read(name))
    except Exception:
        return None


def part_counts(names: list[str]) -> dict[str, int]:
    return {
        "slides": len([n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)]),
        "slide_layouts": len(
            [n for n in names if re.match(r"ppt/slideLayouts/slideLayout\d+\.xml$", n)]
        ),
        "slide_masters": len(
            [n for n in names if re.match(r"ppt/slideMasters/slideMaster\d+\.xml$", n)]
        ),
        "notes": len([n for n in names if re.match(r"ppt/notesSlides/notesSlide\d+\.xml$", n)]),
        "themes": len([n for n in names if n.startswith("ppt/theme/") and n.endswith(".xml")]),
        "media": len([n for n in names if n.startswith("ppt/media/")]),
    }


def parse_theme(zf: zipfile.ZipFile, names: list[str]) -> dict[str, Any]:
    theme_names = sorted(n for n in names if n.startswith("ppt/theme/") and n.endswith(".xml"))
    if not theme_names:
        return {}
    root = xml_root(zf, theme_names[0])
    if root is None:
        return {"path": theme_names[0], "error": "failed to parse"}

    colors: dict[str, str] = {}
    clr_scheme = root.find(".//a:clrScheme", NS)
    if clr_scheme is not None:
        for child in clr_scheme:
            key = etree.QName(child).localname
            srgb = child.find(".//a:srgbClr", NS)
            sysclr = child.find(".//a:sysClr", NS)
            if srgb is not None and srgb.get("val"):
                colors[key] = srgb.get("val").upper()
            elif sysclr is not None and sysclr.get("lastClr"):
                colors[key] = sysclr.get("lastClr").upper()

    fonts: dict[str, dict[str, str]] = {}
    for scheme_name in ("majorFont", "minorFont"):
        scheme = root.find(f".//a:{scheme_name}", NS)
        if scheme is None:
            continue
        entry: dict[str, str] = {}
        for tag in ("latin", "ea", "cs"):
            node = scheme.find(f"a:{tag}", NS)
            if node is not None and node.get("typeface"):
                entry[tag] = node.get("typeface")
        fonts[scheme_name] = entry

    return {"path": theme_names[0], "colors": colors, "fonts": fonts}


def xml_usage(zf: zipfile.ZipFile, names: list[str]) -> dict[str, Any]:
    inspect_prefixes = ("ppt/slides/", "ppt/slideLayouts/", "ppt/slideMasters/")
    xml_names = [
        n
        for n in names
        if n.endswith(".xml") and any(n.startswith(prefix) for prefix in inspect_prefixes)
    ]
    fonts: Counter[str] = Counter()
    srgb_colors: Counter[str] = Counter()
    scheme_colors: Counter[str] = Counter()

    for name in xml_names:
        root = xml_root(zf, name)
        if root is None:
            continue
        for node in root.xpath(".//*[@typeface]"):
            val = node.get("typeface")
            if val:
                fonts[val] += 1
        for node in root.findall(".//a:srgbClr", NS):
            val = node.get("val")
            if val:
                srgb_colors[val.upper()] += 1
        for node in root.findall(".//a:schemeClr", NS):
            val = node.get("val")
            if val:
                scheme_colors[val] += 1

    return {
        "fonts": dict(fonts.most_common(20)),
        "srgb_colors": dict(srgb_colors.most_common(30)),
        "scheme_colors": dict(scheme_colors.most_common(20)),
    }


def media_summary(zf: zipfile.ZipFile, names: list[str]) -> dict[str, Any]:
    media = [n for n in names if n.startswith("ppt/media/")]
    by_ext: defaultdict[str, dict[str, int]] = defaultdict(lambda: {"count": 0, "bytes": 0})
    largest: list[dict[str, Any]] = []
    for name in media:
        info = zf.getinfo(name)
        ext = Path(name).suffix.lower() or "(none)"
        by_ext[ext]["count"] += 1
        by_ext[ext]["bytes"] += info.file_size
        largest.append({"path": name, "bytes": info.file_size})
    largest = sorted(largest, key=lambda x: x["bytes"], reverse=True)[:8]
    return {"by_extension": dict(sorted(by_ext.items())), "largest": largest}


def placeholder_entry(shape: Any) -> dict[str, Any]:
    ph = shape.placeholder_format
    return {
        "idx": ph.idx,
        "type": str(ph.type),
        "left": emu_to_in(shape.left),
        "top": emu_to_in(shape.top),
        "width": emu_to_in(shape.width),
        "height": emu_to_in(shape.height),
    }


def shape_type_name(shape: Any) -> str:
    return str(shape.shape_type).split(" ")[0]


def parse_with_python_pptx(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    slide_w = emu_to_in(prs.slide_width)
    slide_h = emu_to_in(prs.slide_height)

    layout_usage: Counter[str] = Counter()
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name or f"layout-{slide.slide_layout.slide_layout_id}"
        layout_usage[layout_name] += 1
        type_counts: Counter[str] = Counter()
        placeholders: Counter[str] = Counter()
        text_chars = 0
        text_shapes = 0
        pictures = 0
        title = ""

        if slide.shapes.title is not None and getattr(slide.shapes.title, "text", None):
            title = clean_text(slide.shapes.title.text)

        for shape in slide.shapes:
            type_counts[shape_type_name(shape)] += 1
            if shape.is_placeholder:
                placeholders[str(shape.placeholder_format.type)] += 1
            if shape.has_text_frame:
                text = clean_text(shape.text, 5000)
                if text:
                    text_shapes += 1
                    text_chars += len(text)
            if "PICTURE" in shape_type_name(shape):
                pictures += 1

        slides.append(
            {
                "index": idx,
                "layout": layout_name,
                "title": title,
                "shape_types": dict(type_counts),
                "placeholder_types": dict(placeholders),
                "text_shapes": text_shapes,
                "text_chars": text_chars,
                "pictures": pictures,
            }
        )

    layouts: list[dict[str, Any]] = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.shapes:
            if shape.is_placeholder:
                placeholders.append(placeholder_entry(shape))
        layouts.append(
            {
                "index": idx,
                "name": layout.name or f"layout-{idx}",
                "used_by_slides": layout_usage.get(layout.name or f"layout-{idx}", 0),
                "placeholders": placeholders,
            }
        )

    return {
        "slide_size": {"width_in": slide_w, "height_in": slide_h, "aspect": round(slide_w / slide_h, 3)},
        "slide_count": len(prs.slides),
        "layouts": layouts,
        "layout_usage": dict(layout_usage),
        "slides": slides,
    }


def inspect_pptx(path: Path, root: Path) -> dict[str, Any]:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        zip_info = {
            "parts": part_counts(names),
            "theme": parse_theme(zf, names),
            "xml_usage": xml_usage(zf, names),
            "media": media_summary(zf, names),
        }
    pptx_info = parse_with_python_pptx(path)
    stat = path.stat()
    variant_match = re.search(r"(蓝|红|橙|绿|黑)", path.name)
    return {
        "file": relpath(path, root),
        "name": path.stem,
        "variant": variant_match.group(1) if variant_match else "",
        "bytes": stat.st_size,
        "size_mb": round(stat.st_size / 1024 / 1024, 2),
        **zip_info,
        **pptx_info,
    }


def find_templates(root: Path) -> list[Path]:
    skip_parts = {".codex", "generated", "outputs", ".git", "__pycache__"}
    removed_source_parts = {"sysu-" + "medical-ai"}
    return sorted(
        [
            p
            for p in root.rglob("*.pptx")
            if not skip_parts.intersection(p.parts)
            and not removed_source_parts.intersection(p.parts)
        ],
        key=lambda p: (str(p.parent), p.name),
    )


def write_markdown(data: dict[str, Any], out_file: Path) -> None:
    lines: list[str] = []
    lines.append("# SYSU PPT Template Inventory")
    lines.append("")
    lines.append(f"Generated from `{data['root']}`.")
    lines.append("")
    lines.append("## Template Index")
    lines.append("")
    lines.append("| File | Variant | Size MB | Slides | Aspect | Layouts | Theme Fonts | Top Colors |")
    lines.append("|---|---:|---:|---:|---:|---:|---|---|")
    for item in data["templates"]:
        theme_fonts = item.get("theme", {}).get("fonts", {})
        font_bits = []
        for key in ("majorFont", "minorFont"):
            latin = theme_fonts.get(key, {}).get("latin", "")
            ea = theme_fonts.get(key, {}).get("ea", "")
            if latin or ea:
                font_bits.append(f"{key}: {latin or '-'} / {ea or '-'}")
        top_colors = ", ".join(list(item.get("theme", {}).get("colors", {}).values())[:6])
        lines.append(
            "| {file} | {variant} | {size_mb} | {slides} | {aspect} | {layouts} | {fonts} | {colors} |".format(
                file=item["file"],
                variant=item["variant"] or "-",
                size_mb=item["size_mb"],
                slides=item["slide_count"],
                aspect=item["slide_size"]["aspect"],
                layouts=len(item["layouts"]),
                fonts="<br>".join(font_bits) if font_bits else "-",
                colors=top_colors or "-",
            )
        )

    lines.append("")
    lines.append("## Layout Catalog")
    for item in data["templates"]:
        lines.append("")
        lines.append(f"### {item['name']}")
        lines.append("")
        lines.append("| Layout | Used | Placeholder Types | Geometry Summary |")
        lines.append("|---|---:|---|---|")
        for layout in item["layouts"]:
            types = Counter(ph["type"] for ph in layout["placeholders"])
            type_text = ", ".join(f"{k} x{v}" for k, v in types.items()) or "-"
            geoms = []
            for ph in layout["placeholders"][:6]:
                geoms.append(
                    f"{ph['type']}@({ph['left']},{ph['top']},{ph['width']}x{ph['height']})"
                )
            if len(layout["placeholders"]) > 6:
                geoms.append(f"+{len(layout['placeholders']) - 6} more")
            lines.append(
                f"| {layout['index']}: {layout['name']} | {layout['used_by_slides']} | {type_text} | {'; '.join(geoms) or '-'} |"
            )

    lines.append("")
    lines.append("## Slide Inventory")
    for item in data["templates"]:
        lines.append("")
        lines.append(f"### {item['name']}")
        lines.append("")
        lines.append("| Slide | Layout | Title/Text Signal | Text Chars | Pictures |")
        lines.append("|---:|---|---|---:|---:|")
        for slide in item["slides"]:
            title = slide["title"] or "-"
            lines.append(
                f"| {slide['index']} | {slide['layout']} | {title} | {slide['text_chars']} | {slide['pictures']} |"
            )

    lines.append("")
    lines.append("## Extraction Notes")
    lines.append("")
    lines.append("- Slide indices are zero-based when mapping template slides.")
    lines.append("- This inventory intentionally covers source templates under `templates/source/`; generated PPTX templates are routed through `templates/styles/style-index.json`.")
    lines.append("- Theme colors come from `ppt/theme/theme*.xml`; runtime colors also include direct `srgbClr` values found in slides, layouts, and masters.")
    lines.append("- The local machine did not need visual rendering for this inventory. If LibreOffice and Poppler are available, add thumbnail review before final deck delivery.")
    out_file.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--root", default=".", help="Project root to scan")
    parser.add_argument("--out-dir", required=True, help="Directory for JSON and Markdown inventory")
    args = parser.parse_args()

    root = Path(args.root).resolve()
    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    templates = find_templates(root)
    data = {
        "root": str(root).replace("\\", "/"),
        "templates": [inspect_pptx(path, root) for path in templates],
    }

    json_file = out_dir / "template-inventory.json"
    md_file = out_dir / "template-inventory.md"
    json_file.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    write_markdown(data, md_file)
    print(f"Wrote {json_file}")
    print(f"Wrote {md_file}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
