#!/usr/bin/env python
"""Generate spacious Beamer-inspired SYSU PPTX templates.

These PPTX files preserve Beamer's academic structure but redraw it for
PowerPoint projection: larger typography, fewer simultaneous elements, wider
margins, and bigger visual regions than a true LaTeX Beamer PDF.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[4]
STYLE_ROOT = ROOT / "templates" / "styles"
OUTPUT_ROOT = ROOT / "outputs" / "style-showcase" / "beamer-inspired"
TEMPLATE_ROOT = ROOT / "templates" / "generated" / "beamer-inspired"
LOGO_ROOT = ROOT / "templates" / "generated" / "beamer-assets"
MANIFEST_PATH = ROOT / ".codex" / "skills" / "sysu-ppt-generation" / "assets" / "template-manifest.json"

SLIDE_W = 13.333
SLIDE_H = 7.5


STYLES: list[dict[str, Any]] = [
    {
        "id": "beamer-sysu-blue",
        "name": "Beamer SYSU Blue",
        "source": "templates/source/sysu-official/\u4e2d\u5c71\u5927\u5b66\u5e7b\u706f\u7247\u6a21\u677f-\u84dd.pptx",
        "asset_manifest": "templates/assets/strict-sysu-official-blue/asset-manifest.json",
        "accent": "0B4F6C",
        "secondary": "3494BA",
        "emphasis": "029E73",
        "warning": "DE8F05",
        "bg": "FFFFFF",
        "surface": "F3F8FB",
        "surface2": "E6F1F6",
        "border": "D8E5EC",
        "text": "1D2733",
        "muted": "667085",
        "font": "思源黑体 CN Medium",
        "font_fallback": "微软雅黑",
        "use_case": "Academic reports, seminars, and technical talks using SYSU official blue.",
    },
    {
        "id": "beamer-sysu-green",
        "name": "Beamer SYSU Green",
        "source": "templates/source/sysu-official/\u4e2d\u5c71\u5927\u5b66\u5e7b\u706f\u7247\u6a21\u677f-\u7eff.pptx",
        "asset_manifest": "templates/assets/strict-sysu-official-green/asset-manifest.json",
        "accent": "2F6F4E",
        "secondary": "73A45D",
        "emphasis": "029E73",
        "warning": "DE8F05",
        "bg": "FFFFFF",
        "surface": "F4F8F1",
        "surface2": "EAF2E4",
        "border": "DCE8DC",
        "text": "1E2A22",
        "muted": "617067",
        "font": "思源宋体 CN Heavy",
        "font_fallback": "微软雅黑",
        "use_case": "Biomedical, life-science, public-health, and sustainability talks using SYSU official green.",
    },
    {
        "id": "beamer-sysu-red",
        "name": "Beamer SYSU Red",
        "source": "templates/source/sysu-official/\u4e2d\u5c71\u5927\u5b66\u5e7b\u706f\u7247\u6a21\u677f-\u7ea2.pptx",
        "asset_manifest": "templates/assets/strict-sysu-official-red/asset-manifest.json",
        "accent": "8F1D2C",
        "secondary": "C83A3A",
        "emphasis": "0173B2",
        "warning": "DE8F05",
        "bg": "FFFFFF",
        "surface": "FBF4F4",
        "surface2": "F3E5E5",
        "border": "ECDADA",
        "text": "2B1B1D",
        "muted": "735A5D",
        "font": "思源宋体 CN Medium",
        "font_fallback": "微软雅黑",
        "use_case": "Formal academic reports, defenses, and official talks using SYSU red.",
    },
]


def rel(path: Path) -> str:
    return str(path.relative_to(ROOT)).replace("\\", "/")


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.9)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size,
    color,
    font,
    *,
    bold=False,
    align=None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullet(slide, text, x, y, w, style, *, size=16):
    add_rect(slide, x, y + 0.12, 0.13, 0.13, style["secondary"])
    return add_text(slide, text, x + 0.28, y, w - 0.28, 0.28, size, style["text"], style["font_fallback"])


def add_line(slide, x1, y1, x2, y2, color, width=1.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x1, y1, x2, y2, color):
    line = add_line(slide, x1, y1, x2, y2, color, 2.1)
    line.line.end_arrowhead = True
    return line


def safe_image(path: Path) -> bool:
    try:
        with Image.open(path) as image:
            image.verify()
        return True
    except Exception:
        return False


def alpha_transparency(path: Path) -> float:
    with Image.open(path) as image:
        rgba = image.convert("RGBA")
    alpha = rgba.getchannel("A")
    transparent = sum(1 for value in alpha.getdata() if value < 10)
    return transparent / max(1, rgba.width * rgba.height)


def add_picture_contain(slide, image_path: Path | None, x: float, y: float, w: float, h: float):
    if not image_path:
        return None
    with Image.open(image_path) as image:
        ratio = image.width / max(image.height, 1)
    box_ratio = w / h
    if ratio > box_ratio:
        pw = w
        ph = w / ratio
        px = x
        py = y + (h - ph) / 2
    else:
        ph = h
        pw = h * ratio
        px = x + (w - pw) / 2
        py = y
    return slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))


def add_picture_crop(slide, image_path: Path | None, x: float, y: float, w: float, h: float):
    if not image_path:
        add_rect(slide, x, y, w, h, "E5E7EB", "D0D5DD", radius=True)
        return None
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    with Image.open(image_path) as image:
        img_ratio = image.width / max(image.height, 1)
    frame_ratio = w / h
    if img_ratio > frame_ratio:
        crop = 1 - frame_ratio / img_ratio
        pic.crop_left = crop / 2
        pic.crop_right = crop / 2
    elif img_ratio < frame_ratio:
        crop = 1 - img_ratio / frame_ratio
        pic.crop_top = crop / 2
        pic.crop_bottom = crop / 2
    return pic


def save_presentation(prs: Presentation, path: Path) -> Path:
    try:
        prs.save(path)
        return path
    except PermissionError:
        fallback = path.with_name(f"{path.stem}-updated{path.suffix}")
        prs.save(fallback)
        print(f"Locked: {path}; wrote {fallback}")
        return fallback


def load_manifest(style: dict[str, Any]) -> dict[str, Any]:
    return json.loads((ROOT / style["asset_manifest"]).read_text(encoding="utf-8"))


def assets_by_category(manifest: dict[str, Any], category: str, limit: int) -> list[dict[str, Any]]:
    allowed = {".png", ".jpg", ".jpeg"}
    rows = [
        item
        for item in manifest["assets"]
        if item["category"] == category
        and item.get("width")
        and item.get("height")
        and (ROOT / item["file"]).suffix.lower() in allowed
    ]
    rows.sort(key=lambda item: (len(item.get("slides", [])), item.get("width", 0) * item.get("height", 0)), reverse=True)
    kept = []
    for item in rows:
        path = ROOT / item["file"]
        if path.exists() and safe_image(path):
            kept.append(item)
        if len(kept) >= limit:
            break
    return kept


def select_wordmark_logo(manifest: dict[str, Any]) -> Path | None:
    allowed = {".png"}
    candidates = []
    for item in manifest["assets"]:
        if item["category"] != "logo-wordmark-emblem":
            continue
        path = ROOT / item["file"]
        if path.suffix.lower() not in allowed or not path.exists() or not safe_image(path):
            continue
        width = item.get("width") or 0
        height = item.get("height") or 0
        if not width or not height:
            continue
        aspect = width / height
        if not 2.6 <= aspect <= 4.2:
            continue
        transparency = alpha_transparency(path)
        if transparency < 0.35:
            continue
        score = width * height + transparency * 500000 - abs(aspect - 3.45) * 100000
        candidates.append((score, path))
    if candidates:
        return max(candidates, key=lambda row: row[0])[1]
    fallback = assets_by_category(manifest, "logo-wordmark-emblem", 1)
    return ROOT / fallback[0]["file"] if fallback else None


def make_logo_variant(source: Path | None, color_hex: str, name: str) -> Path | None:
    if not source:
        return None
    LOGO_ROOT.mkdir(parents=True, exist_ok=True)
    out = LOGO_ROOT / f"{name}-{color_hex.lower()}.png"
    if out.exists():
        return out
    with Image.open(source) as image:
        rgba = image.convert("RGBA")
    r, g, b = (int(color_hex[i : i + 2], 16) for i in (0, 2, 4))
    pixels = [(r, g, b, a) if a > 0 else (r, g, b, 0) for *_, a in rgba.getdata()]
    recolored = Image.new("RGBA", rgba.size)
    recolored.putdata(pixels)
    recolored.save(out)
    return out


def add_header(slide, style: dict[str, Any], title: str, index: int, logo_on_dark: Path | None, section: str):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, style["bg"])
    add_rect(slide, 0, 0, SLIDE_W, 0.78, style["accent"])
    add_text(slide, title, 0.62, 0.18, 8.5, 0.34, 20, "FFFFFF", style["font"], bold=True)
    add_text(slide, section, 9.25, 0.26, 1.25, 0.18, 10.5, "FFFFFF", style["font_fallback"], align=PP_ALIGN.RIGHT)
    add_picture_contain(slide, logo_on_dark, 10.75, 0.13, 1.62, 0.46)
    add_line(slide, 0.68, 6.96, 12.55, 6.96, style["border"], 0.8)
    add_text(slide, "Sun Yat-sen University", 0.68, 7.1, 3.0, 0.18, 9.8, style["muted"], "Arial")
    add_text(slide, f"{index:02d}", 12.05, 7.1, 0.5, 0.18, 9.8, style["muted"], "Arial", align=PP_ALIGN.RIGHT)


def add_block(slide, style: dict[str, Any], x, y, w, h, label, body, *, kind="normal"):
    header_color = {"normal": style["accent"], "example": style["emphasis"], "alert": style["warning"]}[kind]
    add_rect(slide, x, y, w, h, "FFFFFF", style["border"], radius=True)
    add_rect(slide, x, y, w, 0.48, header_color)
    add_text(slide, label, x + 0.2, y + 0.13, w - 0.4, 0.18, 11.5, "FFFFFF", style["font"], bold=True)
    add_text(slide, body, x + 0.28, y + 0.78, w - 0.56, h - 0.95, 15.5, style["text"], style["font_fallback"])


def slide_cover(prs, style, logo_on_light, hero):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, style["bg"])
    add_rect(slide, 0, 0, SLIDE_W, 0.9, style["accent"])
    add_picture_crop(slide, hero, 7.15, 1.22, 5.2, 4.15)
    add_rect(slide, 7.15, 5.38, 5.2, 0.08, style["secondary"])
    add_picture_contain(slide, logo_on_light, 0.82, 1.08, 2.35, 0.78)
    add_text(slide, "面向课堂讲授的中山大学学术模板", 0.82, 2.2, 5.85, 0.9, 29, style["accent"], style["font"], bold=True)
    add_text(slide, "从大纲到可复查 PPTX 的生成示例", 0.86, 3.28, 5.45, 0.36, 18, style["text"], style["font_fallback"], bold=True)
    add_text(slide, "中文长标题 / 两栏讲授 / 表格 / 流程图 / 总结页", 0.86, 3.9, 5.55, 0.28, 13.5, style["muted"], style["font_fallback"])
    add_rect(slide, 0.86, 4.56, 4.65, 0.055, style["secondary"])
    add_text(slide, style["name"], 0.86, 4.85, 4.6, 0.28, 14.5, style["accent"], style["font"], bold=True)
    add_text(slide, "Presenter / Institute / Date", 0.86, 6.78, 4.1, 0.18, 9.5, style["muted"], "Arial")


def slide_agenda(prs, style, logo_on_dark):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "课程主线：从问题到课堂结论", 2, logo_on_dark, "结构")
    items = [
        ("01", "提出问题", "用真实场景说明为什么需要这一节课。"),
        ("02", "建立框架", "把概念、数据、方法放入同一个分析结构。"),
        ("03", "展示证据", "用一张表格或图示支撑核心判断。"),
        ("04", "形成结论", "用可操作的总结页结束课堂。"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 1.38 + i * 1.02
        add_text(slide, num, 1.05, y, 0.66, 0.28, 17, style["secondary"], "Arial", bold=True)
        add_text(slide, title, 1.92, y - 0.04, 2.35, 0.32, 18.5, style["accent"], style["font"], bold=True)
        add_text(slide, desc, 4.45, y, 5.9, 0.27, 15.2, style["text"], style["font_fallback"])
        add_line(slide, 1.92, y + 0.54, 10.7, y + 0.54, style["border"], 0.7)


def slide_blocks(prs, style, logo_on_dark):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "长中文标题也要保持完整、可读和有层级", 3, logo_on_dark, "标题页")
    add_text(slide, "主题示例", 0.94, 1.22, 1.5, 0.22, 12.5, style["secondary"], style["font"], bold=True)
    add_text(
        slide,
        "生成式人工智能如何改变大学课堂中的知识组织、学习反馈与教师评价",
        0.92,
        1.62,
        7.1,
        1.32,
        26,
        style["accent"],
        style["font"],
        bold=True,
    )
    add_text(
        slide,
        "长标题页不依赖压缩字号解决问题，而是通过换行、留白和副标题分层，让学生先读到本页的中心判断。",
        0.98,
        3.25,
        6.9,
        0.72,
        17,
        style["text"],
        style["font_fallback"],
    )
    add_rect(slide, 8.55, 1.45, 3.65, 3.9, "FFFFFF", style["border"], radius=True)
    add_text(slide, "本页检查点", 8.86, 1.82, 2.6, 0.26, 16, style["accent"], style["font"], bold=True)
    add_bullet(slide, "标题可在 5 秒内读完", 8.88, 2.38, 2.7, style, size=14.5)
    add_bullet(slide, "关键词不被边界截断", 8.88, 3.0, 2.7, style, size=14.5)
    add_bullet(slide, "副标题解释课堂任务", 8.88, 3.62, 2.7, style, size=14.5)
    add_rect(slide, 0.98, 4.78, 6.55, 0.08, style["secondary"])
    add_text(slide, "适用：章节导入、课程问题、报告核心论点", 0.98, 5.1, 6.2, 0.24, 13.8, style["muted"], style["font_fallback"])


def slide_visual(prs, style, logo_on_dark, hero):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "两栏讲授页：左侧讲逻辑，右侧放证据", 4, logo_on_dark, "讲授")
    add_rect(slide, 0.86, 1.16, 5.15, 4.95, "FFFFFF", style["border"], radius=True)
    add_text(slide, "教师讲授要点", 1.16, 1.52, 2.5, 0.28, 17, style["accent"], style["font"], bold=True)
    add_text(
        slide,
        "以“学习反馈”为例，先区分可观察行为、过程性数据和最终评价，再讨论 AI 工具能补足哪一段信息。",
        1.16,
        2.1,
        4.28,
        0.92,
        16.2,
        style["text"],
        style["font_fallback"],
    )
    add_bullet(slide, "学生是否理解：看过程，不只看答案", 1.18, 3.38, 4.2, style, size=14.8)
    add_bullet(slide, "教师是否调整：看反馈是否可执行", 1.18, 3.96, 4.2, style, size=14.8)
    add_bullet(slide, "工具是否可靠：看记录能否复查", 1.18, 4.54, 4.2, style, size=14.8)
    add_rect(slide, 6.38, 1.16, 5.92, 4.95, style["surface"], style["border"], radius=True)
    add_picture_crop(slide, hero, 6.68, 1.48, 5.32, 3.38)
    add_text(slide, "右侧保留一个主要视觉区域，承载案例、图表或课堂截图。", 6.7, 5.18, 4.95, 0.28, 12.5, style["muted"], style["font_fallback"])


def slide_table(prs, style, logo_on_dark):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "表格页：比较三种课堂反馈证据", 5, logo_on_dark, "比较")
    add_text(slide, "投影表格只回答一个问题：哪类证据最适合支持教师下一步行动？", 0.94, 1.05, 9.7, 0.32, 16, style["text"], style["font_fallback"])
    x, y = 1.0, 1.78
    col_w = [3.05, 3.1, 4.35]
    headers = ["证据类型", "适合回答的问题", "课堂使用约束"]
    rows = [
        ["课堂提问", "学生是否掌握关键概念", "记录代表性回答即可"],
        ["过程日志", "学习路径在哪一步卡住", "只保留可解释字段"],
        ["阶段测验", "是否达到教学目标", "突出一个改进动作"],
    ]
    total_w = sum(col_w)
    add_line(slide, x, y, x + total_w, y, style["accent"], 2.4)
    cursor = x
    for text, w in zip(headers, col_w):
        add_text(slide, text, cursor + 0.1, y + 0.2, w - 0.2, 0.22, 12.5, style["accent"], style["font"], bold=True)
        cursor += w
    add_line(slide, x, y + 0.62, x + total_w, y + 0.62, style["border"], 1.0)
    for r, row in enumerate(rows):
        yy = y + 0.88 + r * 0.78
        cursor = x
        for c, (text, w) in enumerate(zip(row, col_w)):
            color = style["secondary"] if r == 1 and c == 2 else style["text"]
            add_text(slide, text, cursor + 0.1, yy, w - 0.2, 0.33, 13.5, color, style["font_fallback"], bold=c == 0)
            cursor += w
        add_line(slide, x, yy + 0.56, x + total_w, yy + 0.56, style["border"], 0.45)
    add_line(slide, x, y + 0.88 + len(rows) * 0.78, x + total_w, y + 0.88 + len(rows) * 0.78, style["accent"], 1.9)


def slide_algorithm(prs, style, logo_on_dark):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "流程图页：从大纲到可复查课件", 6, logo_on_dark, "流程")
    nodes = [
        ("写大纲", "outline.md"),
        ("选模板", "style.json"),
        ("映射页面", "mapping.json"),
        ("质量检查", "qa-notes.md"),
    ]
    for i, (title, note) in enumerate(nodes):
        x = 0.95 + i * 3.05
        add_rect(slide, x, 2.15, 2.2, 1.18, "FFFFFF", style["border"], radius=True)
        add_text(slide, title, x + 0.18, 2.45, 1.84, 0.24, 15, style["accent"], style["font"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, note, x + 0.18, 2.92, 1.84, 0.18, 9.8, style["muted"], "Arial", align=PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            add_arrow(slide, x + 2.32, 2.75, x + 2.92, 2.75, style["secondary"])
    add_rect(slide, 1.0, 4.42, 10.9, 0.92, style["surface"], style["border"], radius=True)
    add_text(slide, "关键要求：每次生成 deck 都留下中间文件，后续教师或代理可以复查风格、内容替换和 QA 结论。", 1.32, 4.75, 10.05, 0.24, 15.2, style["text"], style["font_fallback"])


def slide_diagram(prs, style, logo_on_dark):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "总结页：把课堂结论压缩成三个可执行动作", 7, logo_on_dark, "总结")
    cards = [
        ("先选风格", "用蓝/绿/红或 Beamer 系列确定视觉边界。"),
        ("再映射页面", "根据内容形状选择长标题、两栏、表格或流程图。"),
        ("最后做 QA", "检查尺寸、字体、溢出、资产来源和预览图。"),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.92 + i * 3.82
        add_rect(slide, x, 1.72, 3.28, 3.05, "FFFFFF", style["border"], radius=True)
        add_rect(slide, x, 1.72, 3.28, 0.12, style["secondary"])
        add_text(slide, f"0{i + 1}", x + 0.28, 2.08, 0.62, 0.24, 14, style["secondary"], "Arial", bold=True)
        add_text(slide, title, x + 0.9, 2.04, 1.85, 0.28, 17, style["accent"], style["font"], bold=True)
        add_text(slide, body, x + 0.3, 2.9, 2.65, 0.72, 15, style["text"], style["font_fallback"])
    add_text(slide, "适用作课程结尾、报告建议页或生成代理交付前的检查页。", 1.0, 5.45, 8.6, 0.3, 15, style["muted"], style["font_fallback"])


def slide_refs(prs, style, logo_on_dark):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, style, "交付记录与备份页", 8, logo_on_dark, "交付")
    refs = [
        "源模板与资产包：提供中山大学颜色、字体、标识和校区图像。",
        "style.json：记录本次 deck 的风格、字体、调色板和生成规则。",
        "qa-notes.md：记录尺寸、页数、渲染检查和已知限制。",
    ]
    for i, ref in enumerate(refs):
        y = 1.55 + i * 0.86
        add_text(slide, f"[{i + 1}]", 1.0, y, 0.45, 0.22, 12.5, style["secondary"], "Arial", bold=True)
        add_text(slide, ref, 1.62, y - 0.02, 9.5, 0.32, 14.2, style["text"], style["font_fallback"])
    add_rect(slide, 1.0, 4.75, 10.2, 0.86, style["surface"], style["border"], radius=True)
    add_text(slide, "主讲稿保持可读，细节证据放入备份页或 qa-notes。", 1.28, 5.03, 8.4, 0.22, 14.5, style["accent"], style["font"], bold=True)


def make_deck(style: dict[str, Any]) -> tuple[Path, Path]:
    manifest = load_manifest(style)
    photos = assets_by_category(manifest, "campus-photo-or-background", 4)
    cutouts = assets_by_category(manifest, "cutout-building-or-decorative", 2)
    logo_source = select_wordmark_logo(manifest)
    logo_on_dark = make_logo_variant(logo_source, "FFFFFF", f"{style['id']}-wordmark")
    logo_on_light = make_logo_variant(logo_source, style["accent"], f"{style['id']}-wordmark")
    hero_item = (photos + cutouts)[0] if (photos + cutouts) else None
    hero = ROOT / hero_item["file"] if hero_item else None

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_cover(prs, style, logo_on_light, hero)
    slide_agenda(prs, style, logo_on_dark)
    slide_blocks(prs, style, logo_on_dark)
    slide_visual(prs, style, logo_on_dark, hero)
    slide_table(prs, style, logo_on_dark)
    slide_algorithm(prs, style, logo_on_dark)
    slide_diagram(prs, style, logo_on_dark)
    slide_refs(prs, style, logo_on_dark)

    TEMPLATE_ROOT.mkdir(parents=True, exist_ok=True)
    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
    template_path = TEMPLATE_ROOT / f"{style['id']}-template.pptx"
    showcase_path = OUTPUT_ROOT / f"{style['id']}-showcase.pptx"
    saved_template_path = save_presentation(prs, template_path)
    saved_showcase_path = save_presentation(prs, showcase_path)
    return saved_template_path, saved_showcase_path


def write_style(style: dict[str, Any], template_path: Path, showcase_path: Path) -> dict[str, Any]:
    style_dir = STYLE_ROOT / style["id"]
    style_dir.mkdir(parents=True, exist_ok=True)
    spec = {
        "id": style["id"],
        "name": style["name"],
        "group": "beamer-inspired",
        "source": style["source"],
        "asset_manifest": style["asset_manifest"],
        "template_pptx": rel(template_path),
        "demo_pptx": rel(showcase_path),
        "aspect_ratio": "16:9",
        "slide_size_inches": {"width": SLIDE_W, "height": SLIDE_H},
        "reference": {
            "repo": "https://github.com/Noi1r/beamer-skill/tree/main/beamer",
            "local_path": ".codex/reference-skills/beamer-skill/beamer",
            "commit": "b73d48a",
        },
        "palette": {
            key: style[key]
            for key in ["accent", "secondary", "emphasis", "warning", "bg", "surface", "surface2", "border", "text", "muted"]
        },
        "palette_or_colors": {
            key: style[key]
            for key in ["accent", "secondary", "emphasis", "warning", "bg", "surface", "surface2", "border", "text", "muted"]
        },
        "fonts": {
            "primary": style["font"],
            "fallback": style["font_fallback"],
        },
        "generation_status": "ready",
        "rules": [
            "Use Beamer-like structure but redraw for PPTX projection scale.",
            "Keep 16:9 canvas with larger margins and fewer simultaneous elements.",
            "Use slide titles around 20 pt and body text normally at 16 pt or larger.",
            "Support Chinese teaching pages with long titles, two-column lecture layouts, tables, flowcharts, and summary pages.",
            "Use at most two large blocks per slide.",
            "Give visuals or tables enough space to be read from a room.",
            "Use SYSU extracted logos, campus imagery, and source template colors.",
        ],
        "use_case": style["use_case"],
    }
    (style_dir / "style.json").write_text(json.dumps(spec, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    lines = [
        f"# {style['name']}",
        "",
        f"- Style ID: `{style['id']}`",
        "- Group: `beamer-inspired`",
        f"- Source SYSU template: `{style['source']}`",
        f"- PPTX template: `{rel(template_path)}`",
        f"- Demo PPTX: `{rel(showcase_path)}`",
        "- Aspect ratio: `16:9`",
        f"- Reference: `Noi1r/beamer-skill` at `{spec['reference']['commit']}`",
        "",
        "## Use",
        "",
        style["use_case"],
        "",
        "This is a spacious PPTX adaptation of Beamer academic layout discipline. It keeps frame titles, footlines, blocks, tables, algorithms, diagrams, references, and backup logic, but uses larger projection-safe typography and fewer elements per slide.",
        "",
        "## Palette",
        "",
        "| Role | HEX |",
        "|---|---|",
    ]
    for key, value in spec["palette"].items():
        lines.append(f"| {key} | `#{value}` |")
    (style_dir / "style.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    return {
        "id": style["id"],
        "name": style["name"],
        "group": "beamer-inspired",
        "style_spec": rel(style_dir / "style.json"),
        "template_pptx": rel(template_path),
        "demo_pptx": rel(showcase_path),
        "asset_manifest": style["asset_manifest"],
        "source": style["source"],
        "generation_status": "ready",
    }


def update_style_index(entries: list[dict[str, Any]]) -> None:
    payload = json.loads((STYLE_ROOT / "style-index.json").read_text(encoding="utf-8"))
    new_ids = {entry["id"] for entry in entries}
    retained = [entry for entry in payload.get("styles", []) if entry.get("id") not in new_ids]
    payload["styles"] = retained + entries
    (STYLE_ROOT / "style-index.json").write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def update_template_manifest(entries: list[dict[str, Any]]) -> None:
    if not MANIFEST_PATH.exists():
        return
    payload = json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))
    payload["beamer_inspired_templates"] = [
        {
            "id": entry["id"],
            "template": entry["template_pptx"],
            "demo": entry["demo_pptx"],
            "source": entry["source"],
            "asset_manifest": entry["asset_manifest"],
        }
        for entry in entries
    ]
    MANIFEST_PATH.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main() -> int:
    entries = []
    for style in STYLES:
        template_path, showcase_path = make_deck(style)
        entries.append(write_style(style, template_path, showcase_path))
        print(f"Wrote {template_path}")
        print(f"Wrote {showcase_path}")
    update_style_index(entries)
    update_template_manifest(entries)
    print(f"Wrote {STYLE_ROOT / 'style-index.json'}")
    print(f"Wrote {MANIFEST_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
