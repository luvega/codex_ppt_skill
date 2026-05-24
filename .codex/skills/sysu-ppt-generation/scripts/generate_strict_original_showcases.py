#!/usr/bin/env python
"""Extract template visual assets and generate faithful SYSU showcase decks.

The outputs are not copies of the source PPTX files. They are style-system
decks built from extracted fonts, theme colors, logos/marks, campus imagery,
and reusable visual elements found inside the source templates.
"""

from __future__ import annotations

import hashlib
import json
import re
import shutil
import zipfile
from collections import Counter, defaultdict
from io import BytesIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[4]
STYLE_ROOT = ROOT / "templates" / "styles"
ASSET_ROOT = ROOT / "templates" / "assets"
SHOWCASE_ROOT = ROOT / "outputs" / "style-showcase" / "template-elements"
LEGACY_STRICT_ROOT = ROOT / "outputs" / "style-showcase" / "strict-original"
SLIDE_W = 13.333
SLIDE_H = 7.5
SCALE_X = 1.0
SCALE_Y = 1.0


STRICT_STYLES: list[dict[str, Any]] = [
    {
        "id": "strict-sysu-official-blue",
        "name": "Strict SYSU Official Blue",
        "source": "templates/source/sysu-official/中山大学幻灯片模板-蓝.pptx",
        "accent": "0B4F6C",
        "secondary": "3494BA",
        "primary_fonts": ["思源黑体 CN Medium", "思源黑体 CN Heavy", "等线"],
        "use_case": "Exact official SYSU blue style with extracted marks, campus imagery, and source typography.",
    },
    {
        "id": "strict-sysu-official-green",
        "name": "Strict SYSU Official Green",
        "source": "templates/source/sysu-official/中山大学幻灯片模板-绿.pptx",
        "accent": "2F6F4E",
        "secondary": "73A45D",
        "primary_fonts": ["思源宋体 CN Heavy", "思源黑体 CN Regular", "Raleway", "汉仪旗黑-50S"],
        "use_case": "Exact official SYSU green style with extracted source serif titles, marks, and campus imagery.",
    },
    {
        "id": "strict-sysu-official-red",
        "name": "Strict SYSU Official Red",
        "source": "templates/source/sysu-official/中山大学幻灯片模板-红.pptx",
        "accent": "8F1D2C",
        "secondary": "C83A3A",
        "primary_fonts": ["思源宋体 CN Medium", "思源宋体 CN Heavy", "Calibri"],
        "use_case": "Exact official SYSU red style with extracted ceremonial marks, source serif titles, and campus imagery.",
    },
    {
        "id": "strict-sysu-medical-ai",
        "name": "Strict SYSU Medical AI",
        "source": "templates/source/sysu-medical-ai/人工智能导论（医学五年制）简单模板.pptx",
        "accent": "0B5C7D",
        "secondary": "18A0A6",
        "primary_fonts": ["Arial", "黑体", "微软雅黑", "Microsoft YaHei", "KaiTi"],
        "use_case": "Medical AI courseware style with extracted clinical-computing diagrams and teaching imagery.",
    },
]


def rel(path: Path) -> str:
    return str(path.relative_to(ROOT)).replace("\\", "/")


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def sx(value: float) -> float:
    return value * SCALE_X


def sy(value: float) -> float:
    return value * SCALE_Y


def sf(value: float) -> float:
    return value * min(SCALE_X, SCALE_Y)


def source_palette(style: dict[str, Any]) -> dict[str, str]:
    palettes = {
        "strict-sysu-official-blue": {
            "bg": "FFFFFF",
            "surface": "F3F8FB",
            "surface2": "E6F1F6",
            "text": "1D2733",
            "muted": "667085",
            "quiet": "98A2B3",
            "border": "D8E5EC",
        },
        "strict-sysu-official-green": {
            "bg": "FFFFFF",
            "surface": "F4F8F1",
            "surface2": "EAF2E4",
            "text": "1E2A22",
            "muted": "617067",
            "quiet": "94A39A",
            "border": "DCE8DC",
        },
        "strict-sysu-official-red": {
            "bg": "FFFFFF",
            "surface": "FBF4F4",
            "surface2": "F3E5E5",
            "text": "2B1B1D",
            "muted": "735A5D",
            "quiet": "A99799",
            "border": "ECDADA",
        },
        "strict-sysu-medical-ai": {
            "bg": "FBFEFF",
            "surface": "EEF8FA",
            "surface2": "E1F4F6",
            "text": "182A33",
            "muted": "58727D",
            "quiet": "8DA2AA",
            "border": "CEE7EC",
        },
    }
    palette = palettes.get(style["id"], palettes["strict-sysu-official-blue"]).copy()
    palette["accent"] = style["accent"]
    palette["secondary"] = style["secondary"]
    return palette


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def safe_copy_name(part_name: str, digest: str) -> str:
    suffix = Path(part_name).suffix.lower()
    stem = Path(part_name).stem.lower()
    return f"{stem}-{digest[:10]}{suffix}"


def parse_relationships(zf: zipfile.ZipFile) -> dict[str, list[int]]:
    usage: dict[str, list[int]] = defaultdict(list)
    rel_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    rel_re = re.compile(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$")
    for rel_name in zf.namelist():
        match = rel_re.match(rel_name)
        if not match:
            continue
        slide_no = int(match.group(1))
        root = ET.fromstring(zf.read(rel_name))
        for rel in root.findall("r:Relationship", rel_ns):
            target = rel.attrib.get("Target", "")
            if "media/" not in target:
                continue
            media_name = "ppt/media/" + target.split("media/", 1)[1]
            usage[media_name].append(slide_no)
    return usage


def theme_colors(zf: zipfile.ZipFile) -> dict[str, str]:
    theme_name = next((n for n in zf.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")), "")
    if not theme_name:
        return {}
    text = zf.read(theme_name).decode("utf-8", errors="ignore")
    colors: dict[str, str] = {}
    for block in re.findall(r"<a:(dk1|lt1|dk2|lt2|accent[1-6]|hlink|folHlink)>.*?</a:\1>", text, flags=re.S):
        pass
    for role in ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3", "accent4", "accent5", "accent6", "hlink", "folHlink"]:
        role_match = re.search(rf"<a:{role}>.*?</a:{role}>", text, flags=re.S)
        if not role_match:
            continue
        val_match = re.search(r'lastClr="([0-9A-Fa-f]{6})"|val="([0-9A-Fa-f]{6})"', role_match.group(0))
        if val_match:
            colors[role] = (val_match.group(1) or val_match.group(2)).upper()
    return colors


def xml_fonts_and_colors(zf: zipfile.ZipFile) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    fonts: Counter[str] = Counter()
    colors: Counter[str] = Counter()
    for name in zf.namelist():
        if not (name.startswith("ppt/") and name.endswith(".xml")):
            continue
        text = zf.read(name).decode("utf-8", errors="ignore")
        fonts.update(re.findall(r'typeface="([^"]+)"', text))
        fonts.update(re.findall(r'latin="([^"]+)"', text))
        colors.update(c.upper() for c in re.findall(r'val="([0-9A-Fa-f]{6})"', text))
    font_rows = [{"font": key, "count": value} for key, value in fonts.most_common(40)]
    color_rows = [{"color": key, "count": value} for key, value in colors.most_common(40)]
    return font_rows, color_rows


def classify_asset(path: str, image_meta: dict[str, Any] | None, used_on: list[int]) -> str:
    suffix = Path(path).suffix.lower()
    if suffix == ".wdp":
        return "unsupported-wdp-photo"
    if suffix == ".svg":
        return "svg-mark-or-diagram"
    if not image_meta:
        return "other-media"
    w = image_meta["width"]
    h = image_meta["height"]
    mode = image_meta.get("mode", "")
    has_alpha = image_meta.get("has_alpha", False)
    ratio = w / max(h, 1)
    if has_alpha and (ratio >= 2.0 or w <= 900 or h <= 550):
        return "logo-wordmark-emblem"
    if has_alpha and w >= 900 and h >= 650:
        return "cutout-building-or-decorative"
    if not has_alpha and (w >= 1000 or h >= 800):
        return "campus-photo-or-background"
    if suffix in {".gif", ".tiff", ".emf"}:
        return "diagram-or-special-media"
    if used_on and len(used_on) >= 3:
        return "recurring-template-element"
    return "diagram-icon-or-small-photo"


def image_metadata(data: bytes, suffix: str) -> tuple[dict[str, Any] | None, Image.Image | None]:
    if suffix == ".svg":
        return None, None
    try:
        image = Image.open(BytesIO(data))
        image.load()
    except Exception:
        return None, None
    has_alpha = image.mode in {"RGBA", "LA"} or ("transparency" in image.info)
    meta = {
        "width": image.width,
        "height": image.height,
        "mode": image.mode,
        "has_alpha": bool(has_alpha),
    }
    return meta, image


def extract_assets(style: dict[str, Any]) -> dict[str, Any]:
    source = ROOT / style["source"]
    out_dir = ASSET_ROOT / style["id"]
    media_dir = out_dir / "media"
    if out_dir.exists():
        shutil.rmtree(out_dir)
    media_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(source) as zf:
        usage = parse_relationships(zf)
        font_rows, color_rows = xml_fonts_and_colors(zf)
        theme = theme_colors(zf)
        media_parts = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        assets = []
        preview_items = []
        for part in media_parts:
            data = zf.read(part)
            digest = sha256_bytes(data)
            dest = media_dir / safe_copy_name(part, digest)
            dest.write_bytes(data)
            suffix = dest.suffix.lower()
            meta, image = image_metadata(data, suffix)
            used_on = sorted(set(usage.get(part, [])))
            category = classify_asset(part, meta, used_on)
            item = {
                "source_part": part,
                "file": rel(dest),
                "sha256": digest,
                "bytes": len(data),
                "extension": suffix,
                "category": category,
                "slides": used_on,
            }
            if meta:
                item.update(meta)
            assets.append(item)
            if image is not None and category not in {"unsupported-wdp-photo"}:
                preview_items.append((item, image.copy()))

    category_counts = Counter(item["category"] for item in assets)
    contact_sheet = out_dir / "contact-sheet.jpg"
    write_contact_sheet(preview_items, contact_sheet, style)
    manifest = {
        "id": style["id"],
        "name": style["name"],
        "source": style["source"],
        "asset_root": rel(out_dir),
        "contact_sheet": rel(contact_sheet),
        "theme_colors": theme,
        "observed_colors": color_rows,
        "fonts": font_rows,
        "primary_fonts": style["primary_fonts"],
        "assets": assets,
        "category_counts": dict(category_counts),
    }
    (out_dir / "asset-manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return manifest


def write_contact_sheet(items: list[tuple[dict[str, Any], Image.Image]], path: Path, style: dict[str, Any]) -> None:
    thumbs = []
    for item, image in items[:80]:
        image.thumbnail((220, 135))
        canvas = Image.new("RGB", (250, 180), "white")
        x = (250 - image.width) // 2
        y = 14 + (126 - image.height) // 2
        if image.mode == "RGBA":
            canvas.paste(image, (x, y), image)
        else:
            canvas.paste(image.convert("RGB"), (x, y))
        draw = ImageDraw.Draw(canvas)
        label = f"{Path(item['source_part']).name} | {item['category'][:22]}"
        draw.text((10, 148), label, fill=(35, 35, 35))
        thumbs.append(canvas)
    cols = 4
    rows = max(1, (len(thumbs) + cols - 1) // cols)
    sheet = Image.new("RGB", (cols * 250, rows * 180 + 55), "white")
    draw = ImageDraw.Draw(sheet)
    draw.rectangle((0, 0, sheet.width, 55), fill="#" + style["accent"])
    draw.text((18, 18), f"{style['name']} extracted media contact sheet", fill="white")
    for idx, thumb in enumerate(thumbs):
        x = (idx % cols) * 250
        y = 55 + (idx // cols) * 180
        sheet.paste(thumb, (x, y))
    path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(path, quality=92)


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, radius: bool = False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(sx(x)), Inches(sy(y)), Inches(sx(w)), Inches(sy(h)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float, color: str, font: str, bold: bool = False, align=None):
    box = slide.shapes.add_textbox(Inches(sx(x)), Inches(sy(y)), Inches(sx(w)), Inches(sy(h)))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sf(size))
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_picture_contain(slide, image_path: Path, x: float, y: float, w: float, h: float):
    with Image.open(image_path) as image:
        ratio = image.width / max(image.height, 1)
    box_ratio = w / h
    if ratio > box_ratio:
        pw = w
        ph = w / ratio
        px = x
        py = y + (h - ph) / 2
    else:
        ph = h
        pw = h * ratio
        px = x + (w - pw) / 2
        py = y
    slide.shapes.add_picture(str(image_path), Inches(sx(px)), Inches(sy(py)), width=Inches(sx(pw)), height=Inches(sy(ph)))


def select_assets(manifest: dict[str, Any], category: str, limit: int) -> list[dict[str, Any]]:
    items = [
        item
        for item in manifest["assets"]
        if item["category"] == category and item.get("width") and Path(ROOT / item["file"]).suffix.lower() in {".png", ".jpg", ".jpeg", ".gif", ".tiff", ".emf"}
    ]
    items.sort(key=lambda x: (len(x.get("slides", [])), x.get("width", 0) * x.get("height", 0)), reverse=True)
    return items[:limit]


def top_colors(style: dict[str, Any], manifest: dict[str, Any]) -> list[str]:
    colors = [style["accent"], style["secondary"]]
    colors.extend(manifest.get("theme_colors", {}).values())
    colors.extend(row["color"] for row in manifest.get("observed_colors", [])[:12])
    clean = []
    for color in colors:
        c = color.upper().lstrip("#")
        if re.fullmatch(r"[0-9A-F]{6}", c) and c not in clean:
            clean.append(c)
    return clean[:12]


def keep_first_slide_only(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list)[1:]:
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def blank_layout(prs: Presentation):
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def make_showcase_deck(style: dict[str, Any], manifest: dict[str, Any]) -> Path:
    global SCALE_X, SCALE_Y
    demo_path = SHOWCASE_ROOT / f"{style['id']}-elements-showcase.pptx"
    prs = Presentation(ROOT / style["source"])
    keep_first_slide_only(prs)
    SCALE_X = (prs.slide_width / 914400) / SLIDE_W
    SCALE_Y = (prs.slide_height / 914400) / SLIDE_H
    font = style["primary_fonts"][0]
    pal = source_palette(style)
    layout = blank_layout(prs)
    logos = select_assets(manifest, "logo-wordmark-emblem", 8)
    photos = select_assets(manifest, "campus-photo-or-background", 8)
    cutouts = select_assets(manifest, "cutout-building-or-decorative", 4)
    diagrams = select_assets(manifest, "diagram-icon-or-small-photo", 8)

    # Slide 1 is deliberately kept from the source PPTX, preserving the original cover.

    slide = prs.slides.add_slide(layout)
    add_header(slide, style, "Typography is a template element", "Extracted font families must stay attached to the selected SYSU source style.", 2)
    add_rect(slide, 0.82, 1.45, 11.7, 4.9, pal["surface"], pal["border"], radius=True)
    y = 1.55
    for idx, font_name in enumerate(style["primary_fonts"][:5]):
        size = [30, 23, 19, 16, 13][idx]
        add_text(slide, font_name, 1.12, y, 4.9, 0.38, size, pal["secondary"] if idx == 0 else pal["text"], font_name, bold=idx < 2)
        add_text(slide, "中山大学模板字体示例 / SYSU template font sample", 6.15, y + 0.05, 5.6, 0.3, max(11, size - 8), pal["muted"], font_name)
        y += 0.82
    add_rect(slide, 1.1, 6.08, 10.95, 0.44, pal["surface2"], pal["border"], radius=True)
    top_font_text = " / ".join(row["font"] for row in manifest["fonts"][:6])
    add_text(slide, f"Most frequent fonts: {top_font_text}", 1.28, 6.22, 10.4, 0.18, 9.3, pal["muted"], "Arial")

    slide = prs.slides.add_slide(layout)
    add_header(slide, style, "Color is a system, not a stripe", "Use theme colors and observed source colors; avoid generic vertical-bar decoration.", 3)
    colors = top_colors(style, manifest)
    for i, color in enumerate(colors):
        x = 0.9 + (i % 6) * 1.85
        y = 1.55 + (i // 6) * 1.25
        add_rect(slide, x, y, 1.38, 0.68, color, pal["border"])
        add_text(slide, f"#{color}", x, y + 0.79, 1.38, 0.18, 8.8, pal["muted"], "Arial", align=PP_ALIGN.CENTER)
    add_rect(slide, 0.92, 4.65, 10.85, 1.18, pal["surface"], pal["border"], radius=True)
    add_text(slide, "Design rule", 1.18, 4.88, 1.6, 0.22, 11.5, pal["secondary"], font, bold=True)
    add_text(slide, "Official blue, green, red, and medical AI styles must remain separate visual systems. Build charts, callouts, rules, and diagrams from extracted colors before inventing new ones.", 2.55, 4.86, 8.7, 0.42, 12.5, pal["text"], font)

    slide = prs.slides.add_slide(layout)
    add_header(slide, style, "Extracted marks and recurring template elements", "These assets come from ppt/media and are indexed by source slide usage.", 4)
    mark_items = (logos + cutouts + diagrams)[:12]
    for i, item in enumerate(mark_items):
        x = 0.86 + (i % 4) * 3.0
        y = 1.5 + (i // 4) * 1.55
        add_rect(slide, x, y, 2.58, 1.12, pal["surface"], pal["border"], radius=True)
        add_picture_contain(slide, ROOT / item["file"], x + 0.15, y + 0.12, 0.95, 0.72)
        add_text(slide, Path(item["source_part"]).name, x + 1.2, y + 0.18, 1.18, 0.2, 7.8, pal["text"], "Arial")
        add_text(slide, item["category"], x + 1.2, y + 0.48, 1.18, 0.28, 7.2, pal["muted"], "Arial")
        add_text(slide, f"slides: {','.join(map(str, item.get('slides', [])[:4])) or '-'}", x + 1.2, y + 0.79, 1.18, 0.16, 6.8, pal["quiet"], "Arial")

    slide = prs.slides.add_slide(layout)
    add_header(slide, style, "Imagery sets the rhythm", "Use a source-colored page with hero visual first and concise interpretation second.", 5)
    hero = (photos + cutouts)[:1]
    if hero:
        add_rect(slide, 0.9, 1.38, 11.45, 3.28, pal["surface"], pal["border"], radius=True)
        add_picture_contain(slide, ROOT / hero[0]["file"], 1.05, 1.52, 11.15, 2.95)
    for i, item in enumerate((photos + cutouts)[1:4]):
        x = 0.92 + i * 3.78
        y = 5.05
        add_rect(slide, x, y, 3.35, 1.05, pal["surface"], pal["border"], radius=True)
        add_picture_contain(slide, ROOT / item["file"], x + 0.12, y + 0.12, 1.1, 0.72)
        add_text(slide, item["category"], x + 1.38, y + 0.28, 1.65, 0.24, 8.2, pal["muted"], "Arial")

    slide = prs.slides.add_slide(layout)
    add_header(slide, style, "Applied slide pattern", "A faithful new page combines source color, SYSU assets, and effective-agents page rhythm.", 6)
    add_rect(slide, 0.86, 1.58, 4.1, 4.65, pal["surface"], pal["border"], radius=True)
    if logos:
        add_picture_contain(slide, ROOT / logos[0]["file"], 1.15, 1.88, 1.8, 0.52)
    add_text(slide, "行动标题陈述结论，而不是只写主题", 1.15, 2.72, 3.35, 0.7, 20, pal["secondary"], font, bold=True)
    add_text(slide, "正文保留模板字体、校徽和校区素材；版式借鉴工程化卡片与图示优先的节奏。", 1.17, 3.72, 3.2, 1.1, 12.5, pal["text"], font)
    add_rect(slide, 1.15, 5.42, 2.9, 0.04, pal["secondary"])
    if photos:
        add_picture_contain(slide, ROOT / photos[0]["file"], 5.45, 1.58, 6.85, 4.35)
    else:
        add_rect(slide, 5.45, 1.58, 6.85, 4.35, pal["surface2"], pal["border"])
    add_text(slide, "Use as a generation target: keep source cover unchanged, then build content pages from extracted SYSU elements.", 5.5, 6.13, 6.7, 0.28, 11, pal["muted"], font)

    demo_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(demo_path)
    return demo_path


def add_header(slide, style: dict[str, Any], title: str, subtitle: str, index: int):
    pal = source_palette(style)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(pal["bg"])
    font = style["primary_fonts"][0]
    add_text(slide, title, 0.82, 0.46, 9.2, 0.42, 22, pal["secondary"], font, bold=True)
    add_text(slide, subtitle, 0.84, 0.94, 9.1, 0.24, 9.8, pal["muted"], font)
    add_rect(slide, 0.84, 1.28, 10.45, 0.018, pal["border"])
    add_text(slide, f"{index:02d}", 11.95, 0.65, 0.42, 0.18, 8.5, pal["quiet"], "Arial", align=PP_ALIGN.RIGHT)


def write_style_files(style: dict[str, Any], manifest: dict[str, Any], demo_path: Path) -> dict[str, Any]:
    style_dir = STYLE_ROOT / style["id"]
    style_dir.mkdir(parents=True, exist_ok=True)
    source = ROOT / style["source"]
    spec = {
        "id": style["id"],
        "name": style["name"],
        "group": "strict-original",
        "strict_original": True,
        "source": style["source"],
        "demo_pptx": rel(demo_path),
        "asset_manifest": rel(ASSET_ROOT / style["id"] / "asset-manifest.json"),
        "contact_sheet": manifest["contact_sheet"],
        "fonts": {
            "primary": style["primary_fonts"],
            "observed_top": manifest["fonts"][:20],
        },
        "colors": {
            "accent": style["accent"],
            "secondary": style["secondary"],
            "theme": manifest["theme_colors"],
            "observed_top": manifest["observed_colors"][:20],
        },
        "assets": {
            "root": manifest["asset_root"],
            "category_counts": manifest["category_counts"],
        },
        "use_case": style["use_case"],
        "generation_rules": [
            "Keep the first slide of the generated style showcase as the original source PPT cover.",
            "Use the source PPTX as the layout authority, but use the extracted asset manifest to choose logos, marks, photos, and recurring decorative elements.",
            "Preserve source font families, especially the observed Source Han Serif/Sans variants such as 思源宋体 CN Heavy or 思源黑体 CN Heavy when present.",
            "Use extracted theme and observed colors before adding new colors.",
            "Treat official blue, green, and red as separate visual systems; do not only swap accent colors.",
            "Do not add a generic full-height left vertical stripe unless the source layout itself contains that element.",
            "For technical content pages, the allowed outside aesthetic reference is ppt169_building_effective_agents page rhythm: precise cards, hero diagram first, concise interpretation second; keep source-template background colors.",
            "For generated decks, map content to source slide layouts first, then use extracted assets to fill or extend slides.",
        ],
        "aesthetic_reference": {
            "project": "ppt169_building_effective_agents",
            "source": ".codex/reference-skills/ppt-master/examples/ppt169_building_effective_agents/spec_lock.md",
            "adaptation": "Borrow page rhythm and technical-card discipline only; keep source-template backgrounds and replace fonts, colors, icons, and imagery with SYSU template assets.",
        },
        "removed_scope": "Lingnan PDF is intentionally excluded from strict styles because it is not an editable PPTX template in this project.",
    }
    (style_dir / "style.json").write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    lines = [
        f"# {style['name']}",
        "",
        f"- Style ID: `{style['id']}`",
        "- Group: `strict-original`",
        f"- Source: `{style['source']}`",
        f"- Demo PPTX: `{rel(demo_path)}`",
        f"- Asset manifest: `{spec['asset_manifest']}`",
        f"- Contact sheet: `{manifest['contact_sheet']}`",
        f"- Primary fonts: {', '.join(style['primary_fonts'])}",
        "",
        "## Use",
        "",
        style["use_case"],
        "",
        "This strict style is a reusable source-derived style system, not a byte-for-byte source deck copy. The showcase keeps the original cover slide, then demonstrates extracted fonts, colors, marks, campus imagery, and template-specific visual elements.",
        "",
        "## Asset Categories",
        "",
    ]
    for key, value in sorted(manifest["category_counts"].items()):
        lines.append(f"- `{key}`: {value}")
    (style_dir / "style.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    return {
        "id": style["id"],
        "name": style["name"],
        "group": "strict-original",
        "style_spec": rel(style_dir / "style.json"),
        "demo_pptx": rel(demo_path),
        "asset_manifest": spec["asset_manifest"],
        "source": rel(source),
    }


def make_gallery(entries: list[dict[str, Any]]) -> Path:
    global SCALE_X, SCALE_Y
    SCALE_X = 1.0
    SCALE_Y = 1.0
    gallery = SHOWCASE_ROOT / "template-elements-gallery.pptx"
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")
    add_text(slide, "SYSU Template Element Library", 0.72, 0.55, 8.2, 0.5, 27, "1D2733", "Microsoft YaHei", bold=True)
    add_text(slide, "Strict showcases keep the source cover, then demonstrate extracted SYSU fonts, colors, marks, and imagery.", 0.75, 1.13, 9.5, 0.28, 12, "667085", "Microsoft YaHei")
    for i, entry in enumerate(entries):
        style = next(s for s in STRICT_STYLES if s["id"] == entry["id"])
        spec = json.loads((ROOT / entry["style_spec"]).read_text(encoding="utf-8"))
        x = 0.78 + (i % 2) * 6.18
        y = 1.85 + (i // 2) * 1.68
        add_rect(slide, x, y, 5.48, 1.26, "F8FAFC", "D8E5EC", radius=True)
        add_rect(slide, x + 0.32, y + 0.24, 0.52, 0.08, style["accent"])
        add_text(slide, entry["name"], x + 0.38, y + 0.2, 3.75, 0.26, 12.8, style["accent"], style["primary_fonts"][0], bold=True)
        add_text(slide, entry["id"], x + 0.38, y + 0.55, 3.75, 0.2, 8.5, "667085", "Arial")
        counts = spec["assets"]["category_counts"]
        total = sum(counts.values())
        add_text(slide, f"{total} media assets", x + 3.98, y + 0.55, 1.05, 0.2, 8.6, "1D2733", "Arial", align=PP_ALIGN.RIGHT)
    gallery.parent.mkdir(parents=True, exist_ok=True)
    prs.save(gallery)
    return gallery


def update_style_index(entries: list[dict[str, Any]], gallery: Path) -> None:
    index_path = STYLE_ROOT / "style-index.json"
    payload = {
        "styles": entries,
        "template_element_gallery_pptx": rel(gallery),
    }
    index_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def cleanup_removed_scope() -> None:
    paths = [
        STYLE_ROOT / "strict-lingnan",
        LEGACY_STRICT_ROOT,
    ]
    for path in paths:
        if path.exists():
            resolved = path.resolve()
            if not str(resolved).lower().startswith(str(ROOT.resolve()).lower()):
                raise RuntimeError(f"Refusing to remove outside workspace: {resolved}")
            shutil.rmtree(resolved)


def main() -> int:
    STYLE_ROOT.mkdir(parents=True, exist_ok=True)
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    SHOWCASE_ROOT.mkdir(parents=True, exist_ok=True)
    cleanup_removed_scope()
    entries = []
    for style in STRICT_STYLES:
        manifest = extract_assets(style)
        demo_path = make_showcase_deck(style, manifest)
        entries.append(write_style_files(style, manifest, demo_path))
    gallery = make_gallery(entries)
    update_style_index(entries, gallery)
    print(f"Wrote extracted assets under {ASSET_ROOT}")
    print(f"Wrote template element showcase decks under {SHOWCASE_ROOT}")
    print(f"Wrote {gallery}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
