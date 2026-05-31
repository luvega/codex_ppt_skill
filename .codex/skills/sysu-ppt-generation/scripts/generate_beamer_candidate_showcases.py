#!/usr/bin/env python
"""Generate spacious Beamer-derived SYSU candidate showcase decks.

Each candidate keeps the visual idea of a public Beamer theme, but redraws it
for PowerPoint: larger typography, larger assets, wider spacing, and one color
variant per full slide instead of miniature Beamer-density samples.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[4]
STYLE_ROOT = ROOT / "templates" / "styles"
OUTPUT_ROOT = ROOT / "outputs" / "style-showcase" / "beamer-candidates"
MANIFEST_PATH = ROOT / ".codex" / "skills" / "sysu-ppt-generation" / "assets" / "template-manifest.json"
STYLE_INDEX_PATH = STYLE_ROOT / "style-index.json"
LOGO_ROOT = ROOT / "templates" / "generated" / "beamer-assets"

SLIDE_W = 13.333
SLIDE_H = 7.5

SYSU_ZH = "\u4e2d\u5c71\u5927\u5b66"
FONT_SERIF_HEAVY = "\u601d\u6e90\u5b8b\u4f53 CN Heavy"
FONT_SERIF_MEDIUM = "\u601d\u6e90\u5b8b\u4f53 CN Medium"
FONT_SANS_HEAVY = "\u601d\u6e90\u9ed1\u4f53 CN Heavy"
FONT_SANS_MEDIUM = "\u601d\u6e90\u9ed1\u4f53 CN Medium"
FONT_YAHEI = "\u5fae\u8f6f\u96c5\u9ed1"


VARIANT_PALETTES: dict[str, dict[str, str]] = {
    "blue": {
        "label": "SYSU Blue",
        "strict_id": "strict-sysu-official-blue",
        "accent": "0B4F6C",
        "secondary": "3494BA",
        "emphasis": "029E73",
        "warning": "DE8F05",
        "surface": "F3F8FB",
        "surface2": "E6F1F6",
        "border": "D8E5EC",
        "text": "1D2733",
        "muted": "667085",
    },
    "green": {
        "label": "SYSU Green",
        "strict_id": "strict-sysu-official-green",
        "accent": "2F6F4E",
        "secondary": "73A45D",
        "emphasis": "0B8F67",
        "warning": "C98218",
        "surface": "F4F8F1",
        "surface2": "EAF2E4",
        "border": "DCE8DC",
        "text": "1E2A22",
        "muted": "617067",
    },
    "red": {
        "label": "SYSU Red",
        "strict_id": "strict-sysu-official-red",
        "accent": "8F1D2C",
        "secondary": "C83A3A",
        "emphasis": "0173B2",
        "warning": "DE8F05",
        "surface": "FBF4F4",
        "surface2": "F3E5E5",
        "border": "ECDADA",
        "text": "2B1B1D",
        "muted": "735A5D",
    },
}


STYLE_FAMILIES: list[dict[str, Any]] = [
    {
        "id": "simpleplus-sysu-clean",
        "name": "SimplePlus-SYSU Clean",
        "family": "simpleplus",
        "reference_label": "SimplePlus Beamer Theme",
        "reference_urls": [
            "https://github.com/pm25/SimplePlus-BeamerTheme",
            "https://pt.overleaf.com/latex/templates/simpleplus-beamertheme/wfmfjhdcrdfx",
        ],
        "base_variant": "blue",
        "font_primary": FONT_SANS_MEDIUM,
        "font_heading": FONT_SANS_HEAVY,
        "bg": "FFFFFF",
        "neutral": "F8FAFC",
        "use_case": "Clean academic seminars, course reports, and content-first research talks.",
        "rules": [
            "White canvas, light rules, and restrained color.",
            "Small identity elements, large content regions.",
            "One visual or one block group per slide.",
        ],
    },
    {
        "id": "ustc-thu-sysu-institutional",
        "name": "USTC/THU-SYSU Institutional",
        "family": "institutional",
        "reference_label": "USTC and THU Beamer themes",
        "reference_urls": [
            "https://www.overleaf.com/latex/templates/ustc-presentation-slash-beamer-template/rvpmgprgfhmr",
            "https://ctan.org/tex-archive/macros/latex/contrib/beamer-contrib/themes/thubeamer",
        ],
        "base_variant": "red",
        "font_primary": FONT_SERIF_MEDIUM,
        "font_heading": FONT_SERIF_HEAVY,
        "bg": "FBFBF8",
        "neutral": "F4F1EA",
        "use_case": "Formal university reports, defenses, and institution-facing academic presentations.",
        "rules": [
            "Stronger institutional header and frame number.",
            "Serif-heavy Chinese title hierarchy.",
            "Campus imagery and wordmarks remain prominent.",
        ],
    },
    {
        "id": "moloch-sysu-minimal",
        "name": "Moloch-SYSU Minimal",
        "family": "moloch",
        "reference_label": "Moloch / Metropolis lineage",
        "reference_urls": [
            "https://moloch.ink/",
            "https://github.com/jolars/moloch",
        ],
        "base_variant": "blue",
        "font_primary": FONT_SANS_MEDIUM,
        "font_heading": FONT_SANS_HEAVY,
        "bg": "FAFAF7",
        "neutral": "F0F1EE",
        "use_case": "Technical talks that need maximum reading space and minimal visual noise.",
        "rules": [
            "Large type, sparse content, and a single progress bar.",
            "Light background for SYSU compatibility.",
            "No boxed clutter unless the slide needs a formal block.",
        ],
    },
    {
        "id": "sleek-sysu-research",
        "name": "Sleek-SYSU Research",
        "family": "sleek",
        "reference_label": "Sleek Beamer Theme",
        "reference_urls": [
            "https://cs.overleaf.com/latex/templates/sleek-beamer-theme/zzpczkprdbqs",
        ],
        "base_variant": "green",
        "font_primary": FONT_SANS_MEDIUM,
        "font_heading": FONT_SANS_HEAVY,
        "bg": "F6F7F9",
        "neutral": "ECEFF3",
        "use_case": "Algorithm, code, computational method, and data-heavy research presentations.",
        "rules": [
            "Technical panels are large and few.",
            "Use dark-on-light contrast and clear grid alignment.",
            "Reserve color for state, result, and emphasis.",
        ],
    },
    {
        "id": "river-sysu-atelier",
        "name": "River/Beamer Atelier Inspired",
        "family": "river",
        "reference_label": "River / Beamer Atelier inspired",
        "reference_urls": [
            "https://beameratelier.com/",
        ],
        "base_variant": "blue",
        "font_primary": FONT_SERIF_MEDIUM,
        "font_heading": FONT_SERIF_HEAVY,
        "bg": "FBFBFA",
        "neutral": "F1F2EF",
        "use_case": "Long-form lectures, structured research talks, and teaching-first slide decks.",
        "rules": [
            "Visible structure with quiet progress cues.",
            "Large theorem/equation/table environments.",
            "Good for lectures that need recurring navigation.",
        ],
    },
]


def rel(path: Path) -> str:
    return str(path.relative_to(ROOT)).replace("\\", "/")


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def style_index_by_id() -> dict[str, dict[str, Any]]:
    payload = load_json(STYLE_INDEX_PATH)
    return {entry["id"]: entry for entry in payload.get("styles", [])}


def variant_meta(index: dict[str, dict[str, Any]], variant_name: str) -> dict[str, Any]:
    palette = VARIANT_PALETTES[variant_name].copy()
    strict = index[palette["strict_id"]]
    palette["source"] = strict["source"]
    palette["asset_manifest"] = strict["asset_manifest"]
    return palette


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, *, radius: bool = False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.9)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str, width: float = 1.3):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str):
    line = add_line(slide, x1, y1, x2, y2, color, 2.0)
    line.line.end_arrowhead = True
    return line


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: str,
    font: str,
    *,
    bold: bool = False,
    align: PP_ALIGN | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullet(slide, text: str, x: float, y: float, w: float, variant: dict[str, str], family: dict[str, Any], *, size: float = 15.5):
    add_rect(slide, x, y + 0.12, 0.13, 0.13, variant["secondary"])
    add_text(slide, text, x + 0.3, y, w - 0.3, 0.32, size, variant["text"], family["font_primary"])


def safe_image(path: Path) -> bool:
    try:
        with Image.open(path) as image:
            image.verify()
        return True
    except Exception:
        return False


def alpha_transparency(path: Path) -> float:
    with Image.open(path) as image:
        rgba = image.convert("RGBA")
    alpha = rgba.getchannel("A")
    transparent = sum(1 for value in alpha.getdata() if value < 10)
    return transparent / max(1, rgba.width * rgba.height)


def assets_by_category(manifest_path: str, category: str, limit: int) -> list[Path]:
    manifest = load_json(ROOT / manifest_path)
    allowed = {".png", ".jpg", ".jpeg"}
    rows = [
        item
        for item in manifest.get("assets", [])
        if item.get("category") == category
        and item.get("width")
        and item.get("height")
        and (ROOT / item["file"]).suffix.lower() in allowed
    ]
    rows.sort(key=lambda item: (len(item.get("slides", [])), item.get("width", 0) * item.get("height", 0)), reverse=True)
    paths = []
    for item in rows:
        path = ROOT / item["file"]
        if path.exists() and safe_image(path):
            paths.append(path)
        if len(paths) >= limit:
            break
    return paths


def select_wordmark_logo(manifest_path: str) -> Path | None:
    manifest = load_json(ROOT / manifest_path)
    candidates = []
    for item in manifest.get("assets", []):
        if item.get("category") != "logo-wordmark-emblem":
            continue
        path = ROOT / item["file"]
        if path.suffix.lower() != ".png" or not path.exists() or not safe_image(path):
            continue
        width = item.get("width") or 0
        height = item.get("height") or 0
        if not width or not height:
            continue
        aspect = width / height
        if not 2.6 <= aspect <= 4.2:
            continue
        transparency = alpha_transparency(path)
        if transparency < 0.35:
            continue
        score = width * height + transparency * 500000 - abs(aspect - 3.45) * 100000
        candidates.append((score, path))
    if candidates:
        return max(candidates, key=lambda row: row[0])[1]
    fallback = assets_by_category(manifest_path, "logo-wordmark-emblem", 1)
    return fallback[0] if fallback else None


def make_logo_variant(source: Path | None, color_hex: str, name: str) -> Path | None:
    if not source:
        return None
    LOGO_ROOT.mkdir(parents=True, exist_ok=True)
    out = LOGO_ROOT / f"{name}-{color_hex.lower()}.png"
    if out.exists():
        return out
    with Image.open(source) as image:
        rgba = image.convert("RGBA")
    r, g, b = (int(color_hex[i : i + 2], 16) for i in (0, 2, 4))
    pixels = [(r, g, b, a) if a > 0 else (r, g, b, 0) for *_, a in rgba.getdata()]
    recolored = Image.new("RGBA", rgba.size)
    recolored.putdata(pixels)
    recolored.save(out)
    return out


def add_picture_contain(slide, image_path: Path | None, x: float, y: float, w: float, h: float):
    if not image_path:
        return None
    with Image.open(image_path) as image:
        ratio = image.width / max(image.height, 1)
    box_ratio = w / h
    if ratio > box_ratio:
        pw = w
        ph = w / ratio
        px = x
        py = y + (h - ph) / 2
    else:
        ph = h
        pw = h * ratio
        px = x + (w - pw) / 2
        py = y
    return slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))


def add_picture_crop(slide, image_path: Path | None, x: float, y: float, w: float, h: float):
    if not image_path:
        add_rect(slide, x, y, w, h, "E5E7EB", "D0D5DD", radius=True)
        return None
    pic = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    with Image.open(image_path) as image:
        img_ratio = image.width / max(image.height, 1)
    frame_ratio = w / h
    if img_ratio > frame_ratio:
        crop = 1 - frame_ratio / img_ratio
        pic.crop_left = crop / 2
        pic.crop_right = crop / 2
    elif img_ratio < frame_ratio:
        crop = 1 - img_ratio / frame_ratio
        pic.crop_top = crop / 2
        pic.crop_bottom = crop / 2
    return pic


class AssetPack:
    def __init__(self, variants: dict[str, dict[str, Any]]):
        self.logo_light: dict[str, Path | None] = {}
        self.logo_dark: dict[str, Path | None] = {}
        self.photo: dict[str, Path | None] = {}
        for name, variant in variants.items():
            photos = assets_by_category(variant["asset_manifest"], "campus-photo-or-background", 4)
            decors = assets_by_category(variant["asset_manifest"], "cutout-building-or-decorative", 3)
            logo_source = select_wordmark_logo(variant["asset_manifest"])
            self.logo_light[name] = make_logo_variant(logo_source, variant["accent"], f"{name}-candidate-wordmark")
            self.logo_dark[name] = make_logo_variant(logo_source, "FFFFFF", f"{name}-candidate-wordmark")
            self.photo[name] = (photos + decors)[0] if (photos + decors) else None


def add_family_header(slide, family: dict[str, Any], variant: dict[str, str], title: str, slide_no: int, logo_light: Path | None, logo_dark: Path | None):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, family["bg"])
    fam = family["family"]
    if fam == "institutional":
        add_rect(slide, 0, 0, SLIDE_W, 0.9, variant["accent"])
        add_text(slide, title, 0.76, 0.24, 8.3, 0.34, 19, "FFFFFF", family["font_heading"], bold=True)
        add_picture_contain(slide, logo_dark, 10.88, 0.15, 1.58, 0.5)
    elif fam == "sleek":
        add_rect(slide, 0, 0, SLIDE_W, 0.74, "20262E")
        add_rect(slide, 0, 0.74, SLIDE_W, 0.07, variant["accent"])
        add_text(slide, title, 0.76, 0.22, 8.55, 0.28, 17.5, "FFFFFF", family["font_heading"], bold=True)
        add_picture_contain(slide, logo_dark, 10.98, 0.14, 1.48, 0.46)
    elif fam == "moloch":
        add_text(slide, title, 0.82, 0.55, 9.2, 0.52, 26, variant["accent"], family["font_heading"], bold=True)
        add_line(slide, 0.84, 1.22, 4.1, 1.22, variant["secondary"], 2.0)
        add_picture_contain(slide, logo_light, 10.95, 0.46, 1.52, 0.44)
    elif fam == "river":
        add_line(slide, 0.76, 0.72, 12.25, 0.72, variant["border"], 1.2)
        add_text(slide, title, 0.82, 0.32, 8.5, 0.28, 18.5, variant["text"], family["font_heading"], bold=True)
        add_picture_contain(slide, logo_light, 10.94, 0.2, 1.52, 0.42)
    else:
        add_rect(slide, 0, 0, SLIDE_W, 0.08, variant["accent"])
        add_text(slide, title, 0.78, 0.48, 8.6, 0.36, 21, variant["text"], family["font_heading"], bold=True)
        add_picture_contain(slide, logo_light, 10.95, 0.36, 1.52, 0.44)

    if fam == "moloch":
        add_rect(slide, 0.82, 7.1, 11.55, 0.04, "D8DADC")
        add_rect(slide, 0.82, 7.1, 11.55 * slide_no / 7, 0.04, variant["accent"])
    elif fam == "river":
        for i in range(7):
            color = variant["accent"] if i < slide_no else "D7DBD5"
            add_rect(slide, 0.82 + i * 0.42, 7.08, 0.32, 0.04, color)
    else:
        add_line(slide, 0.78, 7.04, 12.45, 7.04, variant["border"], 0.8)
    add_text(slide, f"{SYSU_ZH} / {family['name']}", 0.82, 7.15, 4.9, 0.18, 9.8, variant["muted"], family["font_primary"])
    add_text(slide, f"{slide_no:02d}", 12.0, 7.15, 0.42, 0.18, 9.8, variant["muted"], "Arial", align=PP_ALIGN.RIGHT)


def add_cover(prs, family: dict[str, Any], variants: dict[str, dict[str, Any]], assets: AssetPack):
    variant = variants[family["base_variant"]]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, family["bg"])
    if family["family"] in {"institutional", "sleek"}:
        top = variant["accent"] if family["family"] == "institutional" else "20262E"
        add_rect(slide, 0, 0, SLIDE_W, 0.96, top)
        add_rect(slide, 0, 0.96, SLIDE_W, 0.08, variant["accent"])
    add_picture_crop(slide, assets.photo[family["base_variant"]], 7.0, 1.18, 5.25, 4.25)
    add_rect(slide, 7.0, 5.44, 5.25, 0.08, variant["secondary"])
    add_picture_contain(slide, assets.logo_light[family["base_variant"]], 0.86, 1.08, 2.35, 0.78)
    add_text(slide, family["name"], 0.86, 2.2, 5.78, 0.72, 30, variant["accent"], family["font_heading"], bold=True)
    add_text(slide, "Spacious PPTX candidate showcase", 0.9, 3.13, 5.45, 0.36, 18.5, variant["text"], family["font_primary"], bold=True)
    add_text(slide, "Blue / green / red SYSU variants, redrawn for projection instead of Beamer PDF density.", 0.9, 3.78, 5.55, 0.48, 13.5, variant["muted"], "Arial")
    add_text(slide, family["reference_label"], 0.9, 4.65, 4.8, 0.24, 12.5, variant["secondary"], family["font_primary"], bold=True)
    for i, key in enumerate(["blue", "green", "red"]):
        pal = variants[key]
        x = 0.92 + i * 1.28
        add_rect(slide, x, 5.35, 0.96, 0.24, pal["accent"])
        add_text(slide, key, x, 5.68, 0.96, 0.16, 9.5, variant["muted"], "Arial", align=PP_ALIGN.CENTER)


def add_scale_rule(prs, family: dict[str, Any], variants: dict[str, dict[str, Any]], assets: AssetPack):
    variant = variants[family["base_variant"]]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_family_header(slide, family, variant, "PPTX Scale Replaces Beamer Density", 2, assets.logo_light[family["base_variant"]], assets.logo_dark[family["base_variant"]])
    add_text(slide, "This family keeps Beamer's academic logic, but uses PowerPoint-sized elements.", 0.95, 1.25, 9.9, 0.38, 18, variant["text"], family["font_primary"])
    pairs = [
        ("Title", "20-30 pt"),
        ("Body", "16-19 pt"),
        ("Blocks", "2 large blocks max"),
        ("Visuals", "half-slide or larger"),
    ]
    for i, (label, value) in enumerate(pairs):
        x = 0.95 + i * 3.0
        add_rect(slide, x, 2.2, 2.35, 1.55, "FFFFFF", variant["border"], radius=True)
        add_text(slide, label, x + 0.22, 2.55, 1.7, 0.22, 14, variant["accent"], family["font_heading"], bold=True)
        add_text(slide, value, x + 0.22, 3.05, 1.8, 0.22, 17, variant["text"], family["font_primary"], bold=True)
    add_text(slide, "The previous Beamer-like samples were too close to LaTeX PDF density; these decks use PPTX-first spacing.", 1.0, 4.8, 9.5, 0.55, 17, variant["text"], family["font_primary"])


def add_variant_slide(prs, family: dict[str, Any], variant_key: str, variants: dict[str, dict[str, Any]], assets: AssetPack, slide_no: int):
    variant = variants[variant_key]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    color_label = {"blue": "蓝色", "green": "绿色", "red": "红色"}[variant_key]
    add_family_header(slide, family, variant, f"{color_label}方案：真实中文两栏讲授页", slide_no, assets.logo_light[variant_key], assets.logo_dark[variant_key])
    add_rect(slide, 0.92, 1.2, 5.8, 4.72, "FFFFFF", variant["border"], radius=True)
    add_picture_crop(slide, assets.photo[variant_key], 1.16, 1.44, 5.32, 3.62)
    add_text(slide, "中山大学源模板提取图像或课程案例图", 1.18, 5.25, 4.8, 0.2, 10.5, variant["muted"], family["font_primary"])
    add_text(slide, "课堂页先讲一个判断", 7.18, 1.38, 4.35, 0.44, 23, variant["accent"], family["font_heading"], bold=True)
    add_text(slide, "右侧文字解释学生应看到什么，而不是把教师要说的全部写进页面。视觉区域保持足够大，便于投影检查。", 7.2, 2.08, 4.55, 1.05, 17, variant["text"], family["font_primary"])
    add_bullet(slide, "标题表达结论，不只写主题。", 7.25, 3.55, 4.1, variant, family)
    add_bullet(slide, "图片、表格和流程图各占一页核心空间。", 7.25, 4.15, 4.1, variant, family)
    add_rect(slide, 7.2, 5.25, 3.8, 0.08, variant["secondary"])


def add_component_slide(prs, family: dict[str, Any], variants: dict[str, dict[str, Any]], assets: AssetPack):
    variant = variants["blue"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_family_header(slide, family, variant, "Reusable Components Are Fewer and Larger", 6, assets.logo_light["blue"], assets.logo_dark["blue"])
    if family["family"] == "sleek":
        add_rect(slide, 0.95, 1.25, 5.85, 4.75, "FFFFFF", "D0D5DD", radius=True)
        add_rect(slide, 0.95, 1.25, 5.85, 0.58, "20262E")
        add_text(slide, "Algorithm 1: Projection-safe layout", 1.18, 1.42, 4.8, 0.2, 12.5, "FFFFFF", "Consolas", bold=True)
        lines = [
            "Input: outline, style, assets",
            "1 choose one main frame purpose",
            "2 place a full-size visual or block",
            "3 keep text above 16 pt",
            "4 validate bounds and spacing",
        ]
        for i, line in enumerate(lines):
            color = variant["secondary"] if i == 3 else variant["text"]
            add_text(slide, line, 1.24, 2.12 + i * 0.62, 4.95, 0.22, 13.2, color, "Consolas", bold=i == 3)
        add_large_note(slide, family, variant, 7.35, 1.7, "Technical deck", "Panels are useful, but each panel must be readable as a slide, not as source code in a PDF.")
    elif family["family"] == "moloch":
        add_text(slide, "score = clarity x brand_fidelity", 1.05, 2.05, 6.0, 0.42, 24, variant["accent"], "Consolas", bold=True)
        add_line(slide, 1.08, 2.78, 6.2, 2.78, variant["secondary"], 2.5)
        add_text(slide, "Use whitespace as the layout system. Add boxes only when the content has a formal role.", 1.08, 3.35, 6.1, 0.72, 18, variant["text"], family["font_primary"])
        add_large_note(slide, family, variant, 7.8, 1.75, "Minimal rule", "One equation, one conclusion, and one progress cue are enough for a technical frame.")
    else:
        add_large_note(slide, family, variant, 0.98, 1.55, "Definition", "A reusable PPTX frame is a large, readable content region with a stable brand position.")
        add_large_note(slide, family, variant, 6.8, 1.55, "Example", "Use two blocks at most. Move proof detail, dense tables, and alternatives to backup slides.")
        add_text(slide, "Two-block maximum keeps the template from becoming a compressed Beamer PDF.", 1.05, 4.7, 9.4, 0.46, 18, variant["text"], family["font_primary"])


def add_large_note(slide, family: dict[str, Any], variant: dict[str, str], x: float, y: float, label: str, body: str):
    add_rect(slide, x, y, 4.95, 2.45, "FFFFFF", variant["border"], radius=True)
    add_rect(slide, x, y, 4.95, 0.56, variant["accent"])
    add_text(slide, label, x + 0.26, y + 0.18, 3.9, 0.16, 12, "FFFFFF", family["font_heading"], bold=True)
    add_text(slide, body, x + 0.3, y + 0.88, 4.25, 0.78, 16.2, variant["text"], family["font_primary"])


def add_reference_close(prs, family: dict[str, Any], variants: dict[str, dict[str, Any]], assets: AssetPack):
    variant = variants["red"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_family_header(slide, family, variant, "Selection Notes and Reuse", 7, assets.logo_light["red"], assets.logo_dark["red"])
    add_text(slide, family["name"], 0.95, 1.28, 5.8, 0.44, 23, variant["accent"], family["font_heading"], bold=True)
    add_text(slide, family["use_case"], 0.95, 1.95, 5.8, 0.7, 17, variant["text"], family["font_primary"])
    for i, rule in enumerate(family["rules"]):
        add_bullet(slide, rule, 1.0, 3.08 + i * 0.58, 5.85, variant, family, size=15.2)
    add_rect(slide, 7.35, 1.35, 4.75, 3.85, "FFFFFF", variant["border"], radius=True)
    add_text(slide, "References", 7.68, 1.68, 2.3, 0.24, 15, variant["accent"], family["font_heading"], bold=True)
    for i, url in enumerate(family["reference_urls"]):
        add_text(slide, f"[{i + 1}] {url}", 7.68, 2.2 + i * 0.58, 3.9, 0.24, 10.5, variant["muted"], "Arial")
    add_text(slide, "Use this style ID after choosing the candidate direction.", 7.68, 4.36, 3.7, 0.36, 13.5, variant["text"], family["font_primary"])


def make_deck(family: dict[str, Any], variants: dict[str, dict[str, Any]], assets: AssetPack) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    add_cover(prs, family, variants, assets)
    add_scale_rule(prs, family, variants, assets)
    add_variant_slide(prs, family, "blue", variants, assets, 3)
    add_variant_slide(prs, family, "green", variants, assets, 4)
    add_variant_slide(prs, family, "red", variants, assets, 5)
    add_component_slide(prs, family, variants, assets)
    add_reference_close(prs, family, variants, assets)

    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)
    path = OUTPUT_ROOT / f"{family['id']}-showcase.pptx"
    prs.save(path)
    return path


def write_style(family: dict[str, Any], showcase_path: Path, variants: dict[str, dict[str, Any]]) -> dict[str, Any]:
    style_dir = STYLE_ROOT / family["id"]
    style_dir.mkdir(parents=True, exist_ok=True)
    base = variants[family["base_variant"]]
    spec = {
        "id": family["id"],
        "name": family["name"],
        "group": "beamer-candidates",
        "demo_pptx": rel(showcase_path),
        "source": base["source"],
        "asset_manifest": base["asset_manifest"],
        "aspect_ratio": "16:9",
        "slide_size_inches": {"width": SLIDE_W, "height": SLIDE_H},
        "reference": {
            "label": family["reference_label"],
            "urls": family["reference_urls"],
        },
        "fonts": {
            "heading": family["font_heading"],
            "primary": family["font_primary"],
            "fallback": FONT_YAHEI,
        },
        "palette_variants": {
            key: {
                "source": value["source"],
                "asset_manifest": value["asset_manifest"],
                "accent": value["accent"],
                "secondary": value["secondary"],
                "emphasis": value["emphasis"],
                "warning": value["warning"],
                "surface": value["surface"],
                "surface2": value["surface2"],
                "border": value["border"],
                "text": value["text"],
                "muted": value["muted"],
            }
            for key, value in variants.items()
        },
        "palette_or_colors": {
            key: {
                "source": value["source"],
                "asset_manifest": value["asset_manifest"],
                "accent": value["accent"],
                "secondary": value["secondary"],
                "emphasis": value["emphasis"],
                "warning": value["warning"],
                "surface": value["surface"],
                "surface2": value["surface2"],
                "border": value["border"],
                "text": value["text"],
                "muted": value["muted"],
            }
            for key, value in variants.items()
        },
        "layout_rules": family["rules"],
        "pptx_spacing_rules": [
            "Use one color variant per full slide in selection decks.",
            "Use slide titles around 18-30 pt and body text normally at 16 pt or larger.",
            "Use large visual regions; avoid miniature slide previews.",
            "Keep 16:9 widescreen dimensions.",
        ],
        "rules": family["rules"]
        + [
            "Use one color variant per full slide in selection decks.",
            "Use slide titles around 18-30 pt and body text normally at 16 pt or larger.",
            "Use large visual regions; avoid miniature slide previews.",
            "Keep 16:9 widescreen dimensions.",
        ],
        "use_case": family["use_case"],
        "generation_status": "style_selection_only",
        "notes": [
            "This is a PPTX adaptation for SYSU style selection, not a LaTeX export.",
            "All generated slides use 16:9 widescreen dimensions.",
            "Do not add a default full-height left stripe.",
        ],
    }
    write_json(style_dir / "style.json", spec)

    lines = [
        f"# {family['name']}",
        "",
        f"- Style ID: `{family['id']}`",
        "- Group: `beamer-candidates`",
        f"- Demo PPTX: `{rel(showcase_path)}`",
        "- Aspect ratio: `16:9`",
        f"- Reference: {family['reference_label']}",
        "",
        "## Use",
        "",
        family["use_case"],
        "",
        "## PPTX Spacing Rules",
        "",
    ]
    for rule in spec["pptx_spacing_rules"]:
        lines.append(f"- {rule}")
    lines.extend(["", "## Color Variants", "", "| Variant | Accent | Source style |", "|---|---|---|"])
    for key, value in variants.items():
        lines.append(f"| {key} | `#{value['accent']}` | `{value['strict_id']}` |")
    lines.extend(["", "## References", ""])
    for url in family["reference_urls"]:
        lines.append(f"- {url}")
    (style_dir / "style.md").write_text("\n".join(lines) + "\n", encoding="utf-8")

    return {
        "id": family["id"],
        "name": family["name"],
        "group": "beamer-candidates",
        "style_spec": rel(style_dir / "style.json"),
        "demo_pptx": rel(showcase_path),
        "asset_manifest": base["asset_manifest"],
        "source": base["source"],
        "generation_status": "style_selection_only",
        "status_reason": "Candidate showcase for choosing a Beamer-derived direction; not a default production template.",
    }


def update_style_index(entries: list[dict[str, Any]]) -> None:
    payload = load_json(STYLE_INDEX_PATH)
    new_ids = {entry["id"] for entry in entries}
    retained = [entry for entry in payload.get("styles", []) if entry.get("id") not in new_ids]
    payload["styles"] = retained + entries
    write_json(STYLE_INDEX_PATH, payload)


def update_template_manifest(entries: list[dict[str, Any]]) -> None:
    payload = load_json(MANIFEST_PATH)
    payload["beamer_candidate_showcases"] = [
        {
            "id": entry["id"],
            "demo": entry["demo_pptx"],
            "source": entry["source"],
            "asset_manifest": entry["asset_manifest"],
        }
        for entry in entries
    ]
    write_json(MANIFEST_PATH, payload)


def main() -> int:
    index = style_index_by_id()
    variants = {name: variant_meta(index, name) for name in ["blue", "green", "red"]}
    assets = AssetPack(variants)
    entries = []
    for family in STYLE_FAMILIES:
        showcase_path = make_deck(family, variants, assets)
        entries.append(write_style(family, showcase_path, variants))
        print(f"Wrote {showcase_path}")
    update_style_index(entries)
    update_template_manifest(entries)
    print(f"Wrote {STYLE_INDEX_PATH}")
    print(f"Wrote {MANIFEST_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
